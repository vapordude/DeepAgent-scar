
import os
import json
import zipfile
import pandas as pd
import PyPDF2
from pptx import Presentation
from docx import Document
import io
import base64
from tqdm import tqdm
from typing import Dict, Any, Optional, Tuple
import asyncio
import whisper
import xml.etree.ElementTree as ET


class FileProcessor:
    """File processor supporting multiple file types"""
    
    def __init__(self):
        # Base directory
        self.base_dir = ""
        
        # Cache directory and file
        self.cache_dir = ""
        self.cache_file = ""
        
        # In-memory cache map: file_name -> converted_text
        self._processed_cache = {}

        # Supported extensions
        self.supported_extensions = {
            # Text files
            '.txt': self._process_text_file,
            '.py': self._process_text_file,
            '.md': self._process_text_file,
            
            # Table files
            '.xlsx': self._process_excel_file,
            '.csv': self._process_csv_file,
            
            # Document files
            '.docx': self._process_docx_file,
            '.pdf': self._process_pdf_file,
            '.pptx': self._process_pptx_file,
            
            # Data files
            '.json': self._process_json_file,
            '.jsonl': self._process_jsonl_file,
            '.jsonld': self._process_jsonld_file,
            '.pdb': self._process_pdb_file,
            '.xml': self._process_xml_file,
            
            # Compressed files
            '.zip': self._process_zip_file,
            
            # Audio files (with Whisper transcription)
            '.mp3': self._process_audio_file,
            '.wav': self._process_audio_file,
            '.flac': self._process_audio_file,
            '.m4a': self._process_audio_file,
            '.ogg': self._process_audio_file,
        }

        # Load Whisper model
        try:
            self.whisper_model = whisper.load_model("base")
        except Exception:
            self.whisper_model = None

    def check_cache(self, file_name: str) -> str | None:
        """Return cached converted text if available for file_name."""
        return self._processed_cache.get(file_name)

    def _save_cache(self) -> None:
        """Persist in-memory cache to disk."""
        try:
            with open(self.cache_file, "w", encoding="utf-8") as f:
                json.dump(self._processed_cache, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def set_base_dir(self, new_base_dir: str) -> None:
        """Update base directory and ensure cache paths exist without clearing in-memory cache."""
        if not new_base_dir:
            return
        self.base_dir = new_base_dir
        self.cache_dir = os.path.join(self.base_dir, "cache")
        self.cache_file = os.path.join(self.cache_dir, "processed_files.json")
        os.makedirs(self.cache_dir, exist_ok=True)
        # Attempt to merge on-disk cache at new location
        try:
            if os.path.exists(self.cache_file):
                with open(self.cache_file, "r", encoding="utf-8") as f:
                    on_disk = json.load(f)
                    if isinstance(on_disk, dict):
                        self._processed_cache = {**on_disk, **self._processed_cache}
        except Exception:
            pass

    def _process_text_file(self, file_name: str) -> str:
        """Process text files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return f"Text file content:\n{content}"
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()
                return f"Text file content (latin-1 encoding):\n{content}"
            except Exception as e:
                return f"Error reading text file: {str(e)}"
        except Exception as e:
            return f"Error reading text file: {str(e)}"
    
    def _process_excel_file(self, file_name: str) -> str:
        """Process Excel files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            result = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                result.append(f"\n--- Sheet: {sheet_name} ---")
                result.append(f"Shape: {df.shape}")
                result.append(f"Columns: {list(df.columns)}")
                result.append(f"Data:\n{df.head(60).to_string()}")
                
                if len(df) > 60:
                    result.append(f"\n... and {len(df) - 60} more rows")
            
            return "\n".join(result)
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"
    
    def _process_csv_file(self, file_name: str) -> str:
        """Process CSV files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            df = pd.read_csv(file_path)
            result = [
                f"CSV file content:",
                f"Shape: {df.shape}",
                f"Columns: {list(df.columns)}",
                f"Data:\n{df.head(60).to_string()}"
            ]
            
            if len(df) > 60:
                result.append(f"\n... and {len(df) - 60} more rows")
            
            return "\n".join(result)
        except Exception as e:
            return f"Error reading CSV file: {str(e)}"
    
    def _process_docx_file(self, file_name: str) -> str:
        """Process Word documents"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            doc = Document(file_path)
            result = ["Word document content:"]
            
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    result.append(f"Paragraph {i+1}: {paragraph.text}")
            
            # Process tables
            for i, table in enumerate(doc.tables):
                result.append(f"\n--- Table {i+1} ---")
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    result.append(" | ".join(row_data))
            
            return "\n".join(result)
        except Exception as e:
            return f"Error reading Word document: {str(e)}"
    
    def _process_pdf_file(self, file_name: str) -> str:
        """Process PDF files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            text_content = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text_content.append(f"PDF file with {len(pdf_reader.pages)} pages:")
                
                max_pages = 3
                for page_num, page in enumerate(pdf_reader.pages):
                    if page_num >= max_pages:
                        text_content.append(f"\n... (PDF has more than {max_pages} pages, remaining pages are omitted)")
                        break
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_content.append(f"\n--- Page {page_num + 1} ---")
                        text_content.append(page_text)

            return "\n".join(text_content)
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"
    
    def _process_pptx_file(self, file_name: str) -> str:
        """Process PowerPoint files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            text_content = ["PowerPoint presentation content:"]
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content.append(f"\n--- Slide {slide_num + 1} ---")
                    text_content.extend(slide_text)
            
            return "\n".join(text_content)
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"
    
    def _process_json_file(self, file_name: str) -> str:
        """Process JSON files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return f"JSON file content:\n{json.dumps(data, indent=2, ensure_ascii=False)}"
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"
    
    def _process_jsonl_file(self, file_name: str) -> str:
        """Process JSONL files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            result = ["JSONL file content:"]
            with open(file_path, 'r', encoding='utf-8') as f:
                for i, line in enumerate(f):
                    if line.strip():
                        try:
                            data = json.loads(line.strip())
                            result.append(f"Line {i+1}: {json.dumps(data, ensure_ascii=False)}")
                        except json.JSONDecodeError:
                            result.append(f"Line {i+1}: Invalid JSON - {line.strip()}")
            
            return "\n".join(result)
        except Exception as e:
            return f"Error reading JSONL file: {str(e)}"
    
    def _process_jsonld_file(self, file_name: str) -> str:
        """Process JSON-LD files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return f"JSON-LD file content:\n{json.dumps(data, indent=2, ensure_ascii=False)}"
        except Exception as e:
            return f"Error reading JSON-LD file: {str(e)}"
    
    def _process_pdb_file(self, file_name: str) -> str:
        """Process PDB files (Protein Data Bank)"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            result = ["PDB file content:"]
            with open(file_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()[:60]
                result.extend(lines)

            return "".join(result)
        except Exception as e:
            return f"Error reading PDB file: {str(e)}"
    
    def _process_xml_file(self, file_name: str) -> str:
        """Process XML files"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            result = ["XML file content:"]
            
            # Parse XML file
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # Get root element info
            result.append(f"Root element: {root.tag}")
            if root.attrib:
                result.append(f"Root attributes: {root.attrib}")
            
            # Function to recursively process XML elements
            def process_element(element, level=0, max_depth=5):
                if level > max_depth:
                    return f"{'  ' * level}... (max depth reached)"
                
                indent = "  " * level
                element_info = f"{indent}{element.tag}"
                
                if element.attrib:
                    attrs_str = ", ".join([f'{k}="{v}"' for k, v in element.attrib.items()])
                    element_info += f" [{attrs_str}]"
                
                if element.text and element.text.strip():
                    text = element.text.strip()
                    if len(text) > 100:
                        text = text[:100] + "..."
                    element_info += f": {text}"
                
                result.append(element_info)
                
                # Process child elements
                for child in element:
                    process_element(child, level + 1, max_depth)
            
            # Process XML structure
            process_element(root)
            
            return "\n".join(result)
        except ET.ParseError as e:
            return f"Error parsing XML file: {str(e)}"
        except Exception as e:
            return f"Error reading XML file: {str(e)}"
    
    def _process_zip_file(self, file_name: str) -> str:
        """Process ZIP files: extract all contents directly into base_dir"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            result = ["ZIP file contents:"]
            
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                file_list = zip_ref.namelist()
                result.append(f"Total files: {len(file_list)}")
                result.append("File list:")
                
                # Extract all files directly to base_dir
                zip_ref.extractall(self.base_dir)
                
                for inner_name in file_list[:20]:  # Show first 20 files
                    # file_info = zip_ref.getinfo(inner_name)
                    result.append(f"{inner_name}")
                
                if len(file_list) > 20:
                    result.append(f"  ... and {len(file_list) - 20} more files")
            
            return "\n".join(result)
        except Exception as e:
            return f"Error reading ZIP file: {str(e)}"

    def _process_audio_file(self, file_name: str) -> str:
        """Process audio files using Whisper for transcription"""
        try:
            file_path = os.path.join(self.base_dir, file_name)
            file_size = os.path.getsize(file_path)
            file_basename = os.path.basename(file_path)
            
            # Transcribe audio file
            result = self.whisper_model.transcribe(file_path)
            transcription = result["text"]
            
            # Format output
            output = [
                f"Audio file: {file_basename}",
                f"File size: {file_size} bytes",
                f"Transcription:",
                transcription
            ]
            
            return "\n".join(output)
        except Exception as e:
            return f"Error processing audio file: {str(e)}"

    def process_file(self, file_name: str) -> str:
        """Process a file based on its extension"""
        # Check cache first
        cached = self.check_cache(file_name)
        if cached is not None:
            return cached
        
        file_path = os.path.join(self.base_dir, file_name)
        if not os.path.exists(file_path):
            return f"Error: File not found: {file_name}"
        
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension in self.supported_extensions:
            try:
                converted = self.supported_extensions[file_extension](file_name)
                # Save to cache
                self._processed_cache[file_name] = converted
                self._save_cache()
                return converted
            except Exception as e:
                return f"Error processing {file_extension} file: {str(e)}"
        else:
            return f"Unsupported file type: {file_extension}\nSupported types: {', '.join(self.supported_extensions.keys())}"

    def _iter_all_files_recursive(self):
        """Yield relative file paths under base_dir recursively, excluding cache dir."""
        for root, dirs, files in os.walk(self.base_dir):
            # Exclude cache directory
            dirs[:] = [d for d in dirs if os.path.join(root, d) != self.cache_dir]
            for fname in files:
                abs_path = os.path.join(root, fname)
                rel_path = os.path.relpath(abs_path, self.base_dir)
                yield rel_path

    async def process_local_dir_files(self) -> dict:
        """
        Process all files under base_dir recursively, convert to text, and store into cache file.
        If archives (e.g., zip) produce new files, they will be picked up by subsequent passes in the same call.
        Returns a dict mapping file_name -> converted_text for files processed in this call (not including pre-cached files).
        """
        processed_now: dict[str, str] = {}
        if not os.path.exists(self.base_dir):
            return processed_now

        max_passes = 5  # safety to avoid infinite loops
        for _ in range(max_passes):
            new_items = 0
            file_list = list(self._iter_all_files_recursive())
            for entry in tqdm(file_list):
                if entry.startswith("cache"):  # double guard
                    continue
                # Skip if already cached
                if self.check_cache(entry) is not None:
                    continue
                converted = await process_file_content(self, entry)
                processed_now[entry] = converted
                new_items += 1
            # If no new files were processed in this pass, stop
            if new_items == 0:
                break
        return processed_now


# # Global file processor instance
# file_processor = FileProcessor()


def set_file_processor_base_dir(file_processor: FileProcessor, new_base_dir: str) -> None:
    """Configure the file processor's base directory for GAIA files at runtime."""
    try:
        file_processor.set_base_dir(new_base_dir)
    except Exception:
        pass


async def process_file_content(file_processor: FileProcessor, file_name: str) -> str:
    """
    Asynchronously process file content
    
    Args:
        file_name: Name of the file to process (relative to base_dir)
        
    Returns:
        Processed file content as string
    """
    try:
        # Run file processing in thread pool to avoid blocking
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(None, file_processor.process_file, file_name)
        return result
    except Exception as e:
        return f"Error processing file: {str(e)}"


def get_openai_function_process_file():
    """OpenAI function definition for file processing"""
    return {
        "type": "function",
        "function": {
            "name": "process_file",
            "description": "Process and extract content from various file types including text files, spreadsheets, documents, data files, compressed archives, and audio files with transcription. This tool can help analyze file contents, extract structured data, understand file formats, and transcribe audio content.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_name": {
                        "type": "string",
                        "description": "The file name to process.",
                        "examples": [
                            "example.xlsx",
                            "report.pdf",
                            "data.csv",
                            "audio.mp3",
                            "config.xml"
                        ]
                    }
                },
                "required": ["file_name"]
            }
        }
    }



async def list_files_in_directory(directory_path: str) -> str:
    """
    List files in a directory with their types and sizes
    
    Args:
        directory_path: Path to the directory
        
    Returns:
        Formatted string listing files
    """
    try:
        if not os.path.exists(directory_path):
            return f"Error: Directory not found: {directory_path}"
        
        if not os.path.isdir(directory_path):
            return f"Error: Path is not a directory: {directory_path}"
        
        files = []
        for item in os.listdir(directory_path):
            item_path = os.path.join(directory_path, item)
            if os.path.isfile(item_path):
                file_size = os.path.getsize(item_path)
                file_ext = os.path.splitext(item)[1].lower()
                files.append(f"  {item} ({file_ext}) - {file_size} bytes")
        
        if not files:
            return f"Directory '{directory_path}' is empty"
        
        result = [f"Files in '{directory_path}':", f"Total files: {len(files)}"]
        result.extend(sorted(files))
        
        return "\n".join(result)
    except Exception as e:
        return f"Error listing directory: {str(e)}"


# Test function
async def _test():
    """Test the file processor with sample files"""
    test_files = [
        "366e2f2b-8632-4ef2-81eb-bc3877489217.pdf",
        "3da89939-209c-4086-8520-7eb734e6b4ef.xlsx",
        "7dd30055-0198-452e-8c25-f73dbe27dcb8.pdb",
        "1f975693-876d-457b-a649-393859e79bf3.mp3",
        "example.xml"
    ]
    
    for file_path in test_files:
        if os.path.exists(file_path):
            print(f"\n--- Testing {file_path} ---")
            result = await process_file_content(file_path)
            print(result)

async def process_all_files():
    """Process all files in the base directory"""
    file_processor = FileProcessor()
    file_processor.set_base_dir("./data/GAIA/files")
    processed_now = await file_processor.process_local_dir_files()
    return processed_now


if __name__ == '__main__':
    # asyncio.run(_test())
    asyncio.run(process_all_files())